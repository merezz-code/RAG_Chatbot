import os
import tempfile
import shutil
import pandas as pd
import sys
import re
from pathlib import Path

from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

try:
    import fitz
except ImportError:  # pragma: no cover - optional dependency
    fitz = None

try:
    from docx import Document
except ImportError:  # pragma: no cover - optional dependency
    Document = None

try:
    import openpyxl
except ImportError:  # pragma: no cover - optional dependency
    openpyxl = None

try:
    import pytesseract
    from PIL import Image
except ImportError:  # pragma: no cover - optional dependency
    pytesseract = None
    Image = None

ROOT_DIR = Path(__file__).parent.parent
sys.path.append(str(ROOT_DIR))
from shared.config import get_settings

settings = get_settings()

# ─────────────────────────────────────────────
# Models
# ─────────────────────────────────────────────

class ExtractResponse(BaseModel):
    filename: str
    text: str
    pages: int
    file_type: str
    char_count: int

# ─────────────────────────────────────────────
# App
# ─────────────────────────────────────────────

app = FastAPI(title="Text Extractor Service", version="1.0.2")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ─────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────

def extract_pdf(path: str):
    if fitz is None:
        raise RuntimeError("PyMuPDF n'est pas installé. Installez 'pymupdf' pour extraire les PDF.")

    doc = fitz.open(path)
    text_blocks = []
    for page in doc:
        page_text = page.get_text("text").strip()
        if page_text:
            text_blocks.append(page_text)
    doc.close()
    return "\n\n".join(text_blocks), len(text_blocks)


def extract_docx(path: str):
    if Document is None:
        raise RuntimeError("python-docx n'est pas installé.")

    doc = Document(path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n".join(paragraphs), 1


def extract_doc(path: str):
    try:
        import aspose.words as aw
    except ImportError as exc:
        raise RuntimeError("Aspose.Words n'est pas installé pour les fichiers .doc.") from exc

    doc = aw.Document(path)
    text = doc.to_string(aw.SaveFormat.TEXT)
    return text, 1


def extract_image(path: str):
    if Image is None or pytesseract is None:
        raise RuntimeError("Pillow et pytesseract sont requis pour l'OCR des images.")

    if os.name == "nt":
        default_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        if os.path.exists(default_path):
            pytesseract.pytesseract.tesseract_cmd = default_path

    with Image.open(path) as image:
        text = pytesseract.image_to_string(image, lang="fra+eng")
    return text.strip(), 1


def _row_contains_header_keywords(row) -> bool:
    for val in row.values:
        if pd.isna(val):
            continue
        s = str(val)
        if "GDSSRC" in s or "3RDNAM" in s:
            return True
    return False


def _looks_like_header_row(row) -> bool:
    non_null = [v for v in row.values if pd.notna(v)]
    if len(non_null) < 2:
        return False

    def is_text_like(v) -> bool:
        if not isinstance(v, str):
            return False
        stripped = v.replace('.', '', 1).replace('-', '', 1).strip()
        return not stripped.isdigit()

    text_like_count = sum(1 for v in non_null if is_text_like(v))
    return text_like_count >= 2


def detect_header_row(df_raw: pd.DataFrame) -> tuple[int, bool]:
    for idx, row in df_raw.iterrows():
        if _row_contains_header_keywords(row):
            return idx, True
    for idx, row in df_raw.iterrows():
        if _looks_like_header_row(row):
            return idx, False
    return 0, False


def extract_xlsx(path: str):
    text_blocks = []
    ext = os.path.splitext(path)[1].lower()

    if ext == ".csv":
        encoding_used = "utf-8"
        try:
            df_raw = pd.read_csv(path, header=None, encoding=encoding_used)
        except UnicodeDecodeError:
            encoding_used = "latin-1"
            df_raw = pd.read_csv(path, header=None, encoding=encoding_used)

        if df_raw.empty:
            return "", 1

        df = pd.read_csv(path, skiprows=0, encoding=encoding_used)
        df.columns = [str(c).strip().replace("\n", "") for c in df.columns]

        text_blocks.append("## Feuille : CSV_Data")
        for _, row in df.iterrows():
            row_items = []
            for col_name in df.columns:
                val = row[col_name]
                if "Unnamed:" in str(col_name):
                    continue
                if pd.notna(val) and str(val).strip() != "":
                    row_items.append(f"{str(col_name).upper().strip()}: {str(val).strip().replace(chr(10), ' ')}")
            if row_items:
                text_blocks.append(" | ".join(row_items))

        return "\n\n".join(text_blocks), 1

    if ext in [".xlsx", ".xlsm"] and openpyxl is not None:
        workbook = openpyxl.load_workbook(path, data_only=True, read_only=True)
        try:
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                lines = [f"## Feuille : {sheet_name}"]
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                    if row_values:
                        lines.append(" | ".join(row_values))
                if len(lines) > 1:
                    text_blocks.extend(lines)
            return "\n\n".join(text_blocks), len(workbook.sheetnames)
        finally:
            workbook.close()

    with pd.ExcelFile(path) as excel_file:
        total_sheets = len(excel_file.sheet_names)

        for sheet_name in excel_file.sheet_names:
            df_raw = pd.read_excel(excel_file, sheet_name=sheet_name, header=None)
            if df_raw.empty:
                continue

            header_row_index, has_specific_header = detect_header_row(df_raw)
            df = pd.read_excel(excel_file, sheet_name=sheet_name, skiprows=header_row_index)
            df.columns = [str(c).strip().replace("\n", "") for c in df.columns]
            df.dropna(how='all', inplace=True)

            if has_specific_header and 'GDSSRC' in df.columns:
                df = df[df['GDSSRC'] != 'good/service']

            text_blocks.append(f"## Feuille : {sheet_name}")
            for _, row in df.iterrows():
                row_items = []
                for col_name in df.columns:
                    val = row[col_name]
                    if "Unnamed:" in str(col_name):
                        continue
                    if pd.notna(val) and str(val).strip() != "":
                        clean_col = str(col_name).upper().strip()
                        clean_val = str(val).strip().replace("\n", " ")
                        row_items.append(f"{clean_col}: {clean_val}")
                if row_items:
                    text_blocks.append(" | ".join(row_items))

    return "\n\n".join(text_blocks), total_sheets


def extract_pptx(path: str):
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text.strip())
        return "\n".join(text_runs), len(prs.slides)
    except Exception:
        return "Extraction PPTX non configurée ou fichier corrompu.", 1


def extract_txt(path: str):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read(), 1

# ─────────────────────────────────────────────
# Routes
# ─────────────────────────────────────────────

@app.get("/health")
async def health():
    return {"status": "ok", "service": "text_extractor"}


@app.post("/extract", response_model=ExtractResponse)
@app.post("/v1/extract", response_model=ExtractResponse)
async def extract(file: UploadFile = File(...)):

    filename = file.filename
    ext = os.path.splitext(filename)[1].lower()

    supported = [
        ".pdf", ".docx", ".doc", ".xlsx", ".xlsm", ".pptx", ".txt", ".md",
        ".adoc", ".html", ".xml", ".json", ".jsonl", ".yaml", ".yml",
        ".xls", ".csv", ".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"
    ]

    if ext not in supported:
        raise HTTPException(status_code=400, detail=f"Format non supporté : {ext}")

    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        shutil.copyfileobj(file.file, tmp)
        tmp_path = tmp.name

    try:
        if ext == ".pdf":
            text, pages = extract_pdf(tmp_path)
        elif ext == ".docx":
            text, pages = extract_docx(tmp_path)
        elif ext == ".doc":
            text, pages = extract_doc(tmp_path)
        elif ext in [".xlsx", ".xlsm", ".xls", ".csv"]:
            text, pages = extract_xlsx(tmp_path)
        elif ext == ".pptx":
            text, pages = extract_pptx(tmp_path)
        elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"]:
            text, pages = extract_image(tmp_path)
        else:
            text, pages = extract_txt(tmp_path)

        if not text.strip():
            raise HTTPException(status_code=422, detail="Aucun texte extrait.")

        return ExtractResponse(
            filename=filename,
            text=text,
            pages=pages,
            file_type=ext.replace(".", ""),
            char_count=len(text)
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=settings.text_extractor_port)